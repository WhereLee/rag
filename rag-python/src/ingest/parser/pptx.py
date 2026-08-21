"""pptx 解析器：每 slide → 标题/文本框/表格/图片（图片走 VLM），按形状 y 坐标排序。

- 标题：slide.shapes.title 有文本 → heading level 1
- 图片：shape_type == PICTURE → blob → 统一转 PNG → VLM
- 组合形状：递归展开（GROUP）
- 母版继承文本：占位符自动带母版文本（shape.text 已含继承内容）
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import DocumentNode, ParseError, Parser, rows_to_markdown
from .vlm import VLMClient, to_png_bytes

logger = logging.getLogger("rag.pptx")

IMAGE_PROMPT = "你是文档解析器。图片是 PPT 幻灯片中的插图（架构图/图表/照片等）。请输出 JSON：{\"description\": \"2-4 句话描述图片内容与关键信息（含图中数字/标签）\", \"text_in_image\": \"图中出现的文字，按出现顺序换行分隔；没有则为空串\"}"


class PptxParser(Parser):
    """PPT 解析：按页输出节点，meta 带 page。"""

    def __init__(self, vlm: Optional[VLMClient] = None, max_slides: int = 200) -> None:
        self.vlm = vlm
        self.max_slides = max_slides

    def parse(self, path: Path) -> List[DocumentNode]:
        path = Path(path)
        prs = Presentation(path)
        source = path.name
        if len(prs.slides) > self.max_slides:
            raise ParseError(f"{source}: 页数超限（{len(prs.slides)} > {self.max_slides}）")
        nodes: List[DocumentNode] = []
        for pno, slide in enumerate(prs.slides, 1):
            items: List[tuple] = []
            # 标题
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.text.strip():
                items.append((title_shape.top or 0,
                              DocumentNode("heading", title_shape.text.strip(),
                                           {"source": source, "page": pno, "level": 1, "index": len(nodes)})))
            self._walk_shapes(slide.shapes, items, pno, source, nodes)
            items.sort(key=lambda t: (t[0] if t[0] is not None else 0))
            nodes.extend(n for _, n in items)
        return nodes

    def _walk_shapes(self, shapes, items: List[tuple], pno: int, source: str, nodes: List[DocumentNode]) -> None:
        title_shape = shapes.title if hasattr(shapes, "title") else None
        for shape in shapes:
            if title_shape is not None and shape == title_shape:
                continue  # 标题已单独输出，避免重复
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._walk_shapes(shape.shapes, items, pno, source, nodes)
                continue
            top = getattr(shape, "top", None)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                except Exception as e:
                    logger.warning("pptx image blob failed p%d: %s", pno, e)
                    items.append((top, DocumentNode("image", "[图片解析失败]",
                                                    {"source": source, "page": pno, "index": len(nodes)})))
                    continue
                items.append((top, DocumentNode("image", self._describe_image(blob, source, pno),
                                                {"source": source, "page": pno, "index": len(nodes)})))
            elif shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                md = rows_to_markdown(rows)
                if md.strip("| -"):
                    items.append((top, DocumentNode("table", md,
                                                    {"source": source, "page": pno, "index": len(nodes)})))
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    items.append((top, DocumentNode("paragraph", text,
                                                    {"source": source, "page": pno, "index": len(nodes)})))

    def _describe_image(self, blob: bytes, source: str, pno: int) -> str:
        if self.vlm is None:
            from .pdf import _default_vlm
            self.vlm = _default_vlm()
        try:
            png = to_png_bytes(blob)
            obj = self.vlm.chat_image(IMAGE_PROMPT, png, "pptx_img")
            desc = (obj.get("description") or "").strip()
            text_in = (obj.get("text_in_image") or "").strip()
            parts = [desc] + ([text_in] if text_in else [])
            return "；".join(p for p in parts if p) or "[图片无法解析]"
        except Exception as e:
            logger.warning("pptx image vlm failed: %s p%d: %s", source, pno, e)
            return "[图片解析失败]"
