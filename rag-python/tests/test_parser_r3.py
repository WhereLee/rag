"""R3 单元测试：docx / xlsx / pptx 解析（真实文件构造 + FakeVLM）。"""
import io
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Emu, Inches as PInches

from ingest.parser.base import rows_to_markdown
from ingest.parser.docx import DocxParser, _heading_level, iter_block_items
from ingest.parser.pptx import PptxParser
from ingest.parser.xlsx import XlsxParser


class FakeVLM:
    def __init__(self, obj=None, raise_exc=None):
        self.obj = obj or {"description": "架构图：A 服务调用 B 服务", "text_in_image": "A\nB"}
        self.raise_exc = raise_exc
        self.calls = 0

    def chat_image(self, prompt, png_bytes, cache_prefix):
        self.calls += 1
        if self.raise_exc:
            raise self.raise_exc
        return self.obj


def make_png() -> bytes:
    """1x1 红色 PNG（python-docx / python-pptx 都能用）。"""
    import struct
    import zlib

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw = b"\x00\xff\x00\x00"
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


# ---------- rows_to_markdown ----------

class TestRowsToMarkdown:
    def test_basic(self):
        md = rows_to_markdown([["姓名", "年龄"], ["张三", "18"], ["李四", "20"]])
        assert md.splitlines()[0] == "| 姓名 | 年龄 |"
        assert md.splitlines()[1] == "|---|---|"
        assert md.splitlines()[2] == "| 张三 | 18 |"

    def test_ragged_rows_padded(self):
        md = rows_to_markdown([["a", "b"], ["c"]])
        lines = md.splitlines()
        assert "| c |  |" in md or "| c |" in lines[2]


# ---------- docx ----------

class TestDocx:
    def test_heading_table_image_order(self, tmp_path):
        p = tmp_path / "doc.docx"
        doc = Document()
        doc.add_heading("第一章 概述", level=1)
        doc.add_paragraph("这是正文段落。")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "列A"
        table.cell(0, 1).text = "列B"
        table.cell(1, 0).text = "值1"
        table.cell(1, 1).text = "值2"
        doc.add_paragraph()
        doc.add_picture(io.BytesIO(make_png()), width=Inches(1))
        doc.save(p)

        fake = FakeVLM()
        nodes = DocxParser(vlm=fake).parse(p)
        types = [n.type for n in nodes]
        assert types == ["heading", "paragraph", "table", "image"], types
        assert nodes[0].text == "第一章 概述" and nodes[0].meta["level"] == 1
        assert "列A" in nodes[2].text and "值2" in nodes[2].text
        assert "架构图" in nodes[3].text and "A" in nodes[3].text
        assert fake.calls == 1
        assert all(n.meta["source"] == "doc.docx" for n in nodes)

    def test_heading_levels(self):
        doc = Document()
        doc.add_heading("H1", level=1)
        doc.add_heading("H2", level=2)
        paras = [p for p in iter_block_items(doc)]
        assert [_heading_level(p) for p in paras] == [1, 2]

    def test_vlm_failure_placeholder(self, tmp_path):
        p = tmp_path / "img.docx"
        doc = Document()
        doc.add_paragraph("正文")
        doc.add_picture(io.BytesIO(make_png()), width=Inches(1))
        doc.save(p)
        fake = FakeVLM(raise_exc=RuntimeError("boom"))
        nodes = DocxParser(vlm=fake).parse(p)
        img = next(n for n in nodes if n.type == "image")
        assert img.text == "[图片解析失败]"

    def test_empty_doc(self, tmp_path):
        p = tmp_path / "empty.docx"
        doc = Document()
        doc.save(p)
        assert DocxParser(vlm=FakeVLM()).parse(p) == []

    def test_custom_style_outline(self, tmp_path):
        p = tmp_path / "style.docx"
        doc = Document()
        para = doc.add_paragraph("自定义标题样式")
        para._p.get_or_add_pPr().get_or_add_outlineLvl().val = 0  # outline level 0 → heading 1
        doc.save(p)
        nodes = DocxParser(vlm=FakeVLM()).parse(p)
        assert nodes[0].type == "heading" and nodes[0].meta["level"] == 1


# ---------- xlsx ----------

class TestXlsx:
    def test_multi_sheet(self, tmp_path):
        p = tmp_path / "book.xlsx"
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "产品"
        ws1.append(["名称", "价格"])
        ws1.append(["A", "10"])
        ws1.append(["B", "20"])
        ws2 = wb.create_sheet("库存")
        ws2.append(["编号"])
        wb.save(p)

        nodes = XlsxParser().parse(p)
        assert [n.type for n in nodes] == ["table", "table"]
        assert nodes[0].meta["sheet"] == "产品"
        assert "| 名称 | 价格 |" in nodes[0].text
        assert nodes[1].meta["sheet"] == "库存"

    def test_merged_cells(self, tmp_path):
        p = tmp_path / "merged.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["表头A", "表头B"])
        ws.append(["合并值", None])
        ws.merge_cells("A2:B2")
        wb.save(p)
        nodes = XlsxParser().parse(p)
        md = nodes[0].text
        assert "| 合并值 |  |" in md

    def test_formula_cached_value(self, tmp_path):
        p = tmp_path / "formula.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["a", "b", "sum"])
        ws.append([1, 2, "=A2+B2"])
        wb.save(p)
        # data_only 无缓存值 → 公式单元格为 None
        nodes = XlsxParser().parse(p)
        assert "| 1 | 2 |  |" in nodes[0].text

    def test_empty_sheet_skipped(self, tmp_path):
        p = tmp_path / "empty.xlsx"
        wb = Workbook()
        wb.active.title = "空表"
        wb.save(p)
        assert XlsxParser().parse(p) == []


# ---------- pptx ----------

class TestPptx:
    def test_title_text_table_image(self, tmp_path):
        p = tmp_path / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = "季度汇报"
        body = slide.placeholders[1]
        body.text = "正文要点第一行"
        # 表格
        shape = slide.shapes.add_table(2, 2, PInches(1), PInches(3), PInches(4), PInches(1))
        shape.table.cell(0, 0).text = "指标"
        shape.table.cell(0, 1).text = "数值"
        shape.table.cell(1, 0).text = "收入"
        shape.table.cell(1, 1).text = "100万"
        # 图片
        slide.shapes.add_picture(io.BytesIO(make_png()), PInches(1), PInches(5), width=PInches(1))
        prs.save(p)

        fake = FakeVLM()
        nodes = PptxParser(vlm=fake).parse(p)
        types = [n.type for n in nodes]
        assert types[0] == "heading" and nodes[0].text == "季度汇报"
        assert nodes[0].meta["page"] == 1
        assert any(n.type == "table" and "收入" in n.text for n in nodes)
        img = next(n for n in nodes if n.type == "image")
        assert "架构图" in img.text
        assert fake.calls == 1

    def test_title_not_duplicated(self, tmp_path):
        """标题占位符不能重复输出（shapes.title 与遍历 shapes 重叠）。"""
        p = tmp_path / "dup.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "唯一标题"
        prs.save(p)
        nodes = PptxParser(vlm=FakeVLM()).parse(p)
        headings = [n for n in nodes if n.type == "heading"]
        assert len(headings) == 1

    def test_image_failure_placeholder(self, tmp_path):
        p = tmp_path / "fail.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(io.BytesIO(make_png()), PInches(1), PInches(1), width=PInches(1))
        prs.save(p)
        fake = FakeVLM(raise_exc=RuntimeError("boom"))
        nodes = PptxParser(vlm=fake).parse(p)
        assert nodes[0].type == "image" and nodes[0].text == "[图片解析失败]"
