# -*- coding: utf-8 -*-
"""R3 边界测试：空文档 / 图片失败占位 / 自定义样式 / 损坏文件 / 超大 sheet / 母版文本。
用法: python scripts/debug/verify_rag_r3_edge.py（需在项目根目录执行）
"""
import io
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "rag-python" / "src"))

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PInches

from ingest.parser.docx import DocxParser
from ingest.parser.pptx import PptxParser
from ingest.parser.xlsx import XlsxParser

fail = 0


def check(name, cond, detail=""):
    global fail
    if cond:
        print(f"[OK] {name}")
    else:
        print(f"[FAIL] {name} -> {detail}")
        fail += 1


class BoomVLM:
    def chat_image(self, prompt, png_bytes, cache_prefix):
        raise TimeoutError("vlm timeout")


def make_png() -> bytes:
    import struct
    import zlib

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00")) + chunk(b"IEND", b""))


tmp = Path(tempfile.mkdtemp(prefix="rag_r3_e_"))

# 1. docx 空文档 → 空产物
p = tmp / "empty.docx"
Document().save(p)
check("docx 空文档", DocxParser(vlm=BoomVLM()).parse(p) == [])

# 2. docx 图片 VLM 失败 → 占位不中断
p = tmp / "img.docx"
doc = Document()
doc.add_paragraph("正文仍然在。")
doc.add_picture(io.BytesIO(make_png()), width=Inches(1))
doc.save(p)
nodes = DocxParser(vlm=BoomVLM()).parse(p)
check("docx 图片失败占位", any(n.type == "image" and n.text == "[图片解析失败]" for n in nodes), nodes)

# 3. docx 自定义样式 outline level → heading
p = tmp / "custom.docx"
doc = Document()
para = doc.add_paragraph("自定义标题")
para._p.get_or_add_pPr().get_or_add_outlineLvl().val = 1  # outline 1 → heading 2
doc.save(p)
nodes = DocxParser(vlm=BoomVLM()).parse(p)
check("docx 自定义样式", nodes[0].type == "heading" and nodes[0].meta["level"] == 2, nodes)

# 4. xlsx 空 sheet → 跳过
p = tmp / "empty.xlsx"
wb = Workbook()
wb.active.title = "空"
wb.save(p)
check("xlsx 空表", XlsxParser().parse(p) == [])

# 5. xlsx 超大 sheet（2 万行）→ 正常产出不崩
p = tmp / "big.xlsx"
wb = Workbook()
ws = wb.active
ws.append(["id", "value"])
for i in range(20000):
    ws.append([i, f"v{i}"])
wb.save(p)
nodes = XlsxParser().parse(p)
lines = nodes[0].text.splitlines()
check("xlsx 2万行", len(nodes) == 1 and len(lines) == 20002, len(lines))

# 6. pptx 标题不重复（shapes.title 与遍历重叠）
p = tmp / "dup.pptx"
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "唯一标题"
prs.save(p)
nodes = PptxParser(vlm=BoomVLM()).parse(p)
headings = [n for n in nodes if n.type == "heading"]
check("pptx 标题唯一", len(headings) == 1 and headings[0].text == "唯一标题", nodes)

# 7. pptx 图片失败 → 占位
p = tmp / "img.pptx"
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(io.BytesIO(make_png()), PInches(1), PInches(1), width=PInches(1))
prs.save(p)
nodes = PptxParser(vlm=BoomVLM()).parse(p)
check("pptx 图片失败占位", nodes[0].type == "image" and nodes[0].text == "[图片解析失败]", nodes)

# 8. 损坏文件 → 抛异常（不静默产出空）
p = tmp / "bad.docx"
p.write_bytes(b"not a real docx zip")
try:
    DocxParser(vlm=BoomVLM()).parse(p)
    check("损坏 docx 抛错", False, "未抛错")
except Exception as e:
    check("损坏 docx 抛错", True, str(e)[:50])

p = tmp / "bad.xlsx"
p.write_bytes(b"PK\x03\x04 not really")
try:
    XlsxParser().parse(p)
    check("损坏 xlsx 抛错", False, "未抛错")
except Exception:
    check("损坏 xlsx 抛错", True)

# 9. pptx 空演示文稿 → 空产物
p = tmp / "empty.pptx"
prs = Presentation()
prs.save(p)
check("pptx 空文档", PptxParser(vlm=BoomVLM()).parse(p) == [])

print(f"\nR3 边界测试: {'全部通过' if fail == 0 else f'{fail} 项失败'}")
sys.exit(1 if fail else 0)
