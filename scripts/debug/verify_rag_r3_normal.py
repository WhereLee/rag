# -*- coding: utf-8 -*-
"""R3 正常测试：docx（标题/表格/图）/ xlsx（多 sheet/合并/公式缓存）/ pptx（标题/表格/图）。
用法: python scripts/debug/verify_rag_r3_normal.py（需在项目根目录执行）
"""
import io
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "rag-python" / "src"))

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PInches

from ingest.parser.docx import DocxParser
from ingest.parser.pptx import PptxParser
from ingest.parser.xlsx import XlsxParser

fail = 0


def check(name, cond, detail=""):
    global fail
    if cond:
        print(f"[OK] {name}")
    else:
        print(f"[FAIL] {name} -> {detail}")
        fail += 1


class FakeVLM:
    def __init__(self):
        self.calls = 0

    def chat_image(self, prompt, png_bytes, cache_prefix):
        self.calls += 1
        return {"description": "架构示意图：模块A 与 模块B 交互", "text_in_image": "模块A\n模块B"}


def make_png() -> bytes:
    import struct
    import zlib

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00")) + chunk(b"IEND", b""))


tmp = Path(tempfile.mkdtemp(prefix="rag_r3_n_"))
fake = FakeVLM()

# 1. docx：标题 2 级 + 段落 + 表格 + 图片
p = tmp / "报告.docx"
doc = Document()
doc.add_heading("总体架构", level=1)
doc.add_heading("模块说明", level=2)
doc.add_paragraph("这是正文段落，描述模块职责。")
t = doc.add_table(rows=3, cols=2)
for r in range(3):
    t.cell(r, 0).text = f"模块{r + 1}"
    t.cell(r, 1).text = f"职责{r + 1}"
doc.add_picture(io.BytesIO(make_png()), width=Inches(1))
doc.save(p)
nodes = DocxParser(vlm=fake).parse(p)
types = [n.type for n in nodes]
check("docx 节点类型", types == ["heading", "heading", "paragraph", "table", "image"], types)
check("docx 标题层级", nodes[0].meta["level"] == 1 and nodes[1].meta["level"] == 2)
check("docx 表格内容", "模块3" in nodes[3].text and "职责3" in nodes[3].text)
check("docx 图片 VLM", "模块A" in nodes[4].text, nodes[4].text)

# 2. xlsx：双 sheet + 合并单元格 + 公式缓存
p = tmp / "台账.xlsx"
wb = Workbook()
ws = wb.active
ws.title = "预算"
ws.append(["项目", "金额", "备注"])
ws.append(["研发", "100", None])
ws.append(["运维", "50", "含服务器"])
ws.append(["合计", "=B2+B3", None])
ws.merge_cells("A4:C4")
ws2 = wb.create_sheet("成员")
ws2.append(["姓名"])
ws2.append(["张三"])
wb.save(p)
nodes = XlsxParser().parse(p)
check("xlsx sheet 数", len(nodes) == 2, nodes)
check("xlsx 表头", "| 项目 | 金额 | 备注 |" in nodes[0].text)
check("xlsx 合并取左上值", "| 合计 |  |  |" in nodes[0].text)
check("xlsx 公式不残留", "=B2+B3" not in nodes[0].text)  # 合并区域吞掉公式，取左上值
check("xlsx sheet meta", nodes[1].meta["sheet"] == "成员")

# 3. pptx：标题 + 正文 + 表格 + 图片
p = tmp / "汇报.pptx"
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "季度总结"
body = slide.placeholders[1]
body.text = "本季度完成重点工作。"
tb = slide.shapes.add_table(2, 2, PInches(1), PInches(3), PInches(4), PInches(1))
tb.table.cell(0, 0).text = "指标"
tb.table.cell(0, 1).text = "结果"
tb.table.cell(1, 0).text = "完成率"
tb.table.cell(1, 1).text = "95%"
slide.shapes.add_picture(io.BytesIO(make_png()), PInches(1), PInches(5), width=PInches(1))
prs.save(p)
nodes = PptxParser(vlm=fake).parse(p)
types = [n.type for n in nodes]
check("pptx 节点类型", types == ["heading", "paragraph", "table", "image"], types)
check("pptx 标题", nodes[0].text == "季度总结" and nodes[0].meta["page"] == 1)
check("pptx 表格", "完成率" in nodes[2].text)
check("pptx 图片 VLM", "模块B" in nodes[3].text, nodes[3].text)

check("VLM 总调用", fake.calls == 2, fake.calls)  # docx 1 + pptx 1

print(f"\nR3 正常测试: {'全部通过' if fail == 0 else f'{fail} 项失败'}")
sys.exit(1 if fail else 0)
