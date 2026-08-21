# -*- coding: utf-8 -*-
"""R5 边界测试：上限触发 / 超时 kill / 空产物 / 伪装类型 / 损坏文件 / 未知扩展名。
用法: python scripts/debug/verify_rag_r5_edge.py（需在项目根目录执行）
"""
import sys
import tempfile
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "rag-python" / "src"))

import pymupdf
from openpyxl import Workbook
from pptx import Presentation

from ingest.parser.base import ParseError
from ingest.parser.pptx import PptxParser
from ingest.parser.xlsx import XlsxParser
from ingest.pipeline import MAX_IMAGE_BYTES, parse_file

fail = 0


def check(name, cond, detail=""):
    global fail
    if cond:
        print(f"[OK] {name}")
    else:
        print(f"[FAIL] {name} -> {detail}")
        fail += 1


tmp = Path(tempfile.mkdtemp(prefix="rag_r5_e_"))

# 1. 未知扩展名 → failed
p = tmp / "hack.exe"
p.write_bytes(b"MZ")
res = parse_file(p)
check("未知扩展名拒绝", res.status == "failed" and "不支持" in res.error, res.error)

# 2. 伪装类型（exe 内容改名 pdf）→ 魔数不符 failed
p = tmp / "fake.pdf"
p.write_bytes(b"MZ fake exe content")
res = parse_file(p)
check("伪装类型拒绝", res.status == "failed" and "伪装" in res.error, res.error)

# 3. 空 txt → 空产物 failed
p = tmp / "empty.txt"
p.write_bytes(b"")
res = parse_file(p)
check("空产物拒绝", res.status == "failed" and "空" in res.error, res.error)

# 4. 纯 BOM txt → 空产物 failed
p = tmp / "bom.txt"
p.write_bytes(b"\xef\xbb\xbf")
res = parse_file(p)
check("纯 BOM 拒绝", res.status == "failed", f"{res.status}: {res.error}")

# 5. 损坏 pdf → failed（不崩溃）
p = tmp / "corrupt.pdf"
p.write_bytes(b"%PDF-1.4\nbroken")
res = parse_file(p)
check("损坏 pdf 不崩溃", res.status == "failed", f"{res.status}: {res.error}")

# 6. 损坏 docx → failed
p = tmp / "corrupt.docx"
p.write_bytes(b"PK\x03\x04 not a real zip")
res = parse_file(p)
check("损坏 docx 不崩溃", res.status == "failed", f"{res.status}: {res.error}")

# 7. 超时 kill：极小 timeout → failed 且主进程快速返回
p = tmp / "slow.txt"
p.write_text("内容", encoding="utf-8")
t0 = time.time()
res = parse_file(p, timeout=0.05)
elapsed = time.time() - t0
check("超时 kill", res.status == "failed" and "超时" in res.error, f"{res.status}: {res.error}")
check("超时快速返回", elapsed < 10, f"{elapsed:.1f}s")

# 8. 超大图片（>20MB 稀疏文件）→ 拒绝
p = tmp / "big.png"
with open(p, "wb") as f:
    f.write(b"\x89PNG")
    f.seek(MAX_IMAGE_BYTES)
    f.write(b"\x00")
res = parse_file(p)
check("超大图片拒绝", res.status == "failed" and "过大" in res.error, res.error)

# 9. xlsx 单元格上限（解析器级）
p = tmp / "big.xlsx"
wb = Workbook()
ws = wb.active
for i in range(1000):
    ws.append([i] * 10)
wb.save(p)
try:
    XlsxParser(max_cells=5000).parse(p)
    check("xlsx 单元格上限", False, "未抛错")
except ParseError as e:
    check("xlsx 单元格上限", "单元格" in str(e), e)

# 10. pptx 页数上限（解析器级）
p = tmp / "big.pptx"
prs = Presentation()
for _ in range(3):
    prs.slides.add_slide(prs.slide_layouts[6])
prs.save(p)
try:
    PptxParser(max_slides=2).parse(p)
    check("pptx 页数上限", False, "未抛错")
except ParseError as e:
    check("pptx 页数上限", "页数超限" in str(e), e)

# 11. 重复页 PDF（每页相同内容）→ 空产物/重复率 failed
p = tmp / "dup.pdf"
doc = pymupdf.open()
for _ in range(3):
    page = doc.new_page(width=595, height=842)
    try:
        page.insert_text((72, 72), "完全相同的页面内容。", fontname="china-s")
    except Exception:
        page.insert_text((72, 72), "same page content.")
doc.save(p)
doc.close()
res = parse_file(p)
check("重复页拒绝", res.status == "failed" and "重复" in res.error, f"{res.status}: {res.error}")

print(f"\nR5 边界测试: {'全部通过' if fail == 0 else f'{fail} 项失败'}")
sys.exit(1 if fail else 0)
