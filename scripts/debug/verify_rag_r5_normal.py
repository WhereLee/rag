# -*- coding: utf-8 -*-
"""R5 正常测试：全格式 pipeline 走通（判定→魔数→解析→清洗→校验→状态记录）。
用法: python scripts/debug/verify_rag_r5_normal.py（需在项目根目录执行）
"""
import io
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "rag-python" / "src"))

import pymupdf
from docx import Document as WDoc
from docx.shared import Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PInches

from ingest.pipeline import parse_file

fail = 0


def check(name, cond, detail=""):
    global fail
    if cond:
        print(f"[OK] {name}")
    else:
        print(f"[FAIL] {name} -> {detail}")
        fail += 1


tmp = Path(tempfile.mkdtemp(prefix="rag_r5_n_"))


def make_png() -> bytes:
    import struct
    import zlib

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00")) + chunk(b"IEND", b""))


# 1. txt
p = tmp / "笔记.txt"
p.write_text("第一段内容，描述系统。\n\n第二段内容，补充说明。\n\n第三段内容。", encoding="utf-8")
res = parse_file(p)
check("txt 状态", res.status == "success", res.error)
check("txt 节点", len(res.nodes) == 3 and all(n.type == "paragraph" for n in res.nodes))
check("txt 记录完整", res.file == "笔记.txt" and res.duration > 0)

# 2. md
p = tmp / "文档.md"
p.write_text("# 第一章\n\n正文段落内容。\n\n- 列表项一\n- 列表项二", encoding="utf-8")
res = parse_file(p)
check("md 状态", res.status == "success", res.error)
check("md 标题", res.nodes[0].type == "heading" and res.nodes[0].meta["level"] == 1)
check("md 列表", any(n.type == "list" for n in res.nodes), [n.type for n in res.nodes])

# 3. pdf（文本型 + 页眉 + 页脚 + 图片块；图片走真实 VLM 会慢/花钱——此处跳过图片块，
#    用纯文本 PDF 验证文本路径；图片块路径 R2/R4 已覆盖）
p = tmp / "说明.pdf"
doc = pymupdf.open()
for i in range(2):
    page = doc.new_page(width=595, height=842)
    try:
        page.insert_text((72, 72), f"第{i + 1}页的说明文字内容。", fontname="china-s")
        page.insert_text((72, 40), "统一页眉标题", fontsize=9, fontname="china-s")
    except Exception:
        page.insert_text((72, 72), f"Page {i + 1} content text.")
        page.insert_text((72, 40), "header", fontsize=9)
doc.save(p)
doc.close()
res = parse_file(p)
check("pdf 状态", res.status == "success", res.error)
check("pdf 页眉清洗", all("统一页眉标题" not in n.text and "header" not in n.text for n in res.nodes))
check("pdf 页码", all(n.meta.get("page") in (1, 2) for n in res.nodes))

# 4. docx
p = tmp / "报告.docx"
d = WDoc()
d.add_heading("总体架构", level=1)
d.add_paragraph("正文段落描述模块职责。")
d.save(p)
res = parse_file(p)
check("docx 状态", res.status == "success", res.error)
check("docx 标题", res.nodes[0].type == "heading")

# 5. xlsx
p = tmp / "数据.xlsx"
wb = Workbook()
ws = wb.active
ws.append(["指标", "数值"])
ws.append(["收入", "100"])
wb.save(p)
res = parse_file(p)
check("xlsx 状态", res.status == "success", res.error)
check("xlsx 表格", res.nodes[0].type == "table" and "收入" in res.nodes[0].text)

# 6. pptx
p = tmp / "汇报.pptx"
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "季度总结"
slide.placeholders[1].text = "重点内容说明。"
prs.save(p)
res = parse_file(p)
check("pptx 状态", res.status == "success", res.error)
check("pptx 标题", res.nodes[0].type == "heading")

# 7. png 独立图片（真实 VLM 会调 API——此处验证类型判定与入口，VLM 失败降级为 partial/失败占位亦可接受）
p = tmp / "示意图.png"
pix = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 200, 100))
pix.clear_with(200)
p.write_bytes(pix.tobytes("png"))
res = parse_file(p, timeout=60)
# 真实 VLM 可能成功（image 节点）也可能失败（占位 → partial）；绝不能是"子进程输出异常"类 failed
check("png 入口正常", res.status in ("success", "partial"), f"{res.status}: {res.error}")

print(f"\nR5 正常测试: {'全部通过' if fail == 0 else f'{fail} 项失败'}")
sys.exit(1 if fail else 0)
